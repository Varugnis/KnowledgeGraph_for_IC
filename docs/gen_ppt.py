"""
生成课程PPT：半导体知识图谱构建、推理与大模型增强
运行：python docs/gen_ppt.py
输出：docs/课程PPT.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import sys
import os

# ── 颜色主题 ──────────────────────────────────────────────
C_BG        = RGBColor(0x1A, 0x1A, 0x2E)   # 深蓝背景
C_ACCENT    = RGBColor(0x4A, 0x90, 0xD9)   # 蓝色强调
C_ACCENT2   = RGBColor(0xE6, 0x7E, 0x22)   # 橙色强调
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xC8, 0xD8, 0xE8)   # 浅蓝文字
C_CARD      = RGBColor(0x16, 0x21, 0x3E)   # 卡片背景
C_GREEN     = RGBColor(0x27, 0xAE, 0x60)
C_RED       = RGBColor(0xE7, 0x4C, 0x3C)
C_GRAY      = RGBColor(0x7F, 0x8C, 0x8D)
C_YELLOW    = RGBColor(0xF3, 0x9C, 0x12)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank = prs.slide_layouts[6]   # 完全空白
    return prs.slides.add_slide(blank)


def fill_bg(slide, color=C_BG):
    """填充幻灯片背景色"""
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    """添加矩形色块"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "微软雅黑"
    return txBox


def add_multiline(slide, lines, l, t, w, h,
                  size=14, bold_first=False, color=C_WHITE,
                  line_color=None, spacing=1.1):
    """添加多行文本，lines = [(text, bold, color_override), ...]"""
    txBox = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    import lxml.etree as etree

    first = True
    for item in lines:
        if isinstance(item, str):
            text, bold, clr = item, False, color
        else:
            text = item[0]
            bold = item[1] if len(item) > 1 else False
            clr  = item[2] if len(item) > 2 else color

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.space_after = Pt(2)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = clr
        run.font.name = "微软雅黑"
    return txBox


def add_divider(slide, t, color=C_ACCENT):
    add_rect(slide, 0.5, t, 12.3, 0.04, fill=color)


def slide_header(slide, title, subtitle=None):
    """左侧彩色竖条 + 标题"""
    add_rect(slide, 0, 0, 0.12, 7.5, fill=C_ACCENT)
    add_text(slide, title, 0.3, 0.18, 11, 0.6,
             size=28, bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.3, 0.82, 11, 0.4,
                 size=14, color=C_LIGHT)
    add_divider(slide, 1.18)


def card(slide, l, t, w, h, title=None, title_size=14):
    """圆角卡片背景"""
    r = add_rect(slide, l, t, w, h, fill=C_CARD, line=C_ACCENT, line_w=Pt(1))
    if title:
        add_rect(slide, l, t, w, 0.38, fill=C_ACCENT)
        add_text(slide, title, l+0.12, t+0.04, w-0.2, 0.34,
                 size=title_size, bold=True, color=C_WHITE)
    return r


# ═══════════════════════════════════════════════════════════
# SLIDE 1 — 封面
# ═══════════════════════════════════════════════════════════
def slide1(prs):
    s = blank_slide(prs)
    fill_bg(s)

    # 顶部装饰条
    add_rect(s, 0, 0, 13.33, 0.5, fill=C_ACCENT)
    add_rect(s, 0, 7.0, 13.33, 0.5, fill=C_ACCENT)

    # 中央大标题
    add_text(s, "半导体领域知识图谱", 1.2, 1.4, 11, 1.1,
             size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "构建、推理与大模型增强", 1.2, 2.5, 11, 0.8,
             size=30, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    add_divider(s, 3.55)

    # 副标题信息
    add_text(s, "知识图谱课程作业", 1.2, 3.75, 11, 0.5,
             size=18, color=C_LIGHT, align=PP_ALIGN.CENTER)
    add_text(s, "小组成员：成员A  ·  成员B  ·  成员C", 1.2, 4.35, 11, 0.5,
             size=16, color=C_LIGHT, align=PP_ALIGN.CENTER)
    add_text(s, "2026 年 5 月", 1.2, 4.9, 11, 0.5,
             size=16, color=C_GRAY, align=PP_ALIGN.CENTER)

    # 右下角装饰标签
    for i, (label, col) in enumerate([
        ("知识图谱构建", C_ACCENT),
        ("规则 & 路径推理", C_ACCENT2),
        ("GraphRAG 增强", C_GREEN),
    ]):
        add_rect(s, 9.2, 5.4 + i*0.48, 3.6, 0.38, fill=col)
        add_text(s, label, 9.3, 5.44 + i*0.48, 3.4, 0.32,
                 size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 2 — 研究背景与任务目标
# ═══════════════════════════════════════════════════════════
def slide2(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "研究背景与任务目标", "为什么选半导体领域？大模型有何不足？")

    # 左列：选题背景
    card(s, 0.3, 1.35, 6.1, 2.7, "选题背景")
    add_multiline(s, [
        ("半导体是电子信息核心领域，实体关系结构清晰", False, C_WHITE),
        ("全球芯片厂商、工艺节点、材料器件知识密集", False, C_WHITE),
        ("代工关系、工艺继代、竞争格局等适合图谱表示", False, C_WHITE),
        ("工艺型号、代工关系是大模型的高频出错点", False, C_YELLOW),
    ], 0.5, 1.85, 5.8, 2.0, size=13)

    # 左列：大模型不足
    card(s, 0.3, 4.2, 6.1, 2.8, "大模型存在的不足")
    add_multiline(s, [
        ("❌  事实幻觉：工艺型号模糊（4nm左右）、代工厂归属错误", False, C_RED),
        ("❌  缺乏证据：给出结论但无可追溯来源", False, C_RED),
        ("❌  多跳推理弱：跨实体比较逻辑断层", False, C_RED),
        ("❌  知识滞后：无法反映最新工艺动态", False, C_RED),
    ], 0.5, 4.7, 5.8, 2.0, size=13)

    # 右列：任务目标
    card(s, 6.7, 1.35, 6.3, 5.65, "三大任务目标")
    for i, (num, title, desc, col) in enumerate([
        ("01", "图谱构建", "构建半导体知识图谱\n≥110实体、≥318三元组\n8类实体、16类关系", C_ACCENT),
        ("02", "知识推理", "实现规则推理（Horn规则）\n和路径推理（BFS多跳）\n补全隐含知识、发现证据链", C_ACCENT2),
        ("03", "大模型增强", "GraphRAG：图谱子图检索\n注入Prompt，对比有无KG\n的准确性与幻觉差异", C_GREEN),
    ]):
        ty = 1.75 + i * 1.75
        add_rect(s, 6.9, ty, 0.55, 0.55, fill=col)
        add_text(s, num, 6.9, ty, 0.55, 0.55,
                 size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(s, title, 7.55, ty, 2.0, 0.5,
                 size=14, bold=True, color=col)
        add_text(s, desc, 7.55, ty+0.45, 5.2, 1.0,
                 size=12, color=C_LIGHT)


# ═══════════════════════════════════════════════════════════
# SLIDE 3 — 知识图谱 Schema
# ═══════════════════════════════════════════════════════════
def slide3(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "知识图谱 Schema 设计", "8类实体 · 16类关系")

    etype_data = [
        ("Company",      "15", C_ACCENT,   "Intel / TSMC / NVIDIA"),
        ("Chip",         "28", C_ACCENT2,  "H100 / A17 Pro / 骁龙8 Gen 3"),
        ("ProcessNode",  "10", C_GREEN,    "TSMC 3nm / Intel 4 / SMIC N+1"),
        ("Architecture", " 8", RGBColor(0x8E,0x44,0xAD), "Zen4 / Raptor Cove / Cortex-X4"),
        ("Material",     "12", C_RED,      "硅 / 碳化硅 / 氮化镓 / 氧化铪"),
        ("Device",       "12", RGBColor(0x1A,0xBC,0x9C), "FinFET / GAA / HBM / NPU"),
        ("Application",  "15", C_YELLOW,   "AI训练 / 数据中心 / 5G基站"),
        ("Technology",   "10", C_GRAY,     "EUV光刻 / 原子层沉积 / 先进封装"),
    ]

    for i, (name, cnt, col, ex) in enumerate(etype_data):
        row, col_idx = divmod(i, 4)
        lx = 0.3 + col_idx * 3.3
        ty = 1.4 + row * 2.8
        add_rect(s, lx, ty, 3.1, 0.45, fill=col)
        add_text(s, f"{name}  ({cnt}个)", lx+0.1, ty+0.04, 2.9, 0.38,
                 size=13, bold=True, color=C_WHITE)
        add_rect(s, lx, ty+0.45, 3.1, 1.1, fill=C_CARD, line=col, line_w=Pt(1))
        add_text(s, ex, lx+0.1, ty+0.5, 2.9, 1.0, size=11, color=C_LIGHT)

    # 关系列表
    add_rect(s, 0.3, 6.85, 12.7, 0.5, fill=C_CARD)
    rels = "designs · fabricates · uses_process · based_on_arch · used_for · uses_device · uses_material · competes_with · successor_of · enables_process · developed_by · supports_memory · uses_technology · masters_technology · manufactures_tool · requires_device"
    add_text(s, "关系类型（16类）：" + rels, 0.4, 6.88, 12.5, 0.44,
             size=10, color=C_LIGHT)


# ═══════════════════════════════════════════════════════════
# SLIDE 4 — 数据来源与构建流程
# ═══════════════════════════════════════════════════════════
def slide4(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "数据来源与图谱构建流程", "结构化人工整理 → NetworkX 有向多重图")

    # 流程箭头
    steps = [
        ("①\n公开资料", "官网规格\n行业媒体\n技术白皮书", C_ACCENT),
        ("②\nSchema\n设计", "8类实体\n16类关系\n约束规则", C_ACCENT2),
        ("③\n人工\n整理", "entities.csv\ntriples.csv\n置信度标注", C_GREEN),
        ("④\nbuild_kg\n.py", "NetworkX\nMultiDiGraph\n110节点/318边", RGBColor(0x8E,0x44,0xAD)),
        ("⑤\n存储\n输出", "kg_data.json\ntriples.csv\nHTML可视化", C_YELLOW),
    ]
    for i, (step, detail, col) in enumerate(steps):
        lx = 0.4 + i * 2.55
        add_rect(s, lx, 1.4, 2.1, 1.6, fill=col)
        add_text(s, step, lx, 1.4, 2.1, 1.6,
                 size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_rect(s, lx, 3.1, 2.1, 1.1, fill=C_CARD, line=col, line_w=Pt(1))
        add_text(s, detail, lx+0.08, 3.15, 1.95, 1.0, size=11, color=C_LIGHT)
        if i < 4:
            add_text(s, "▶", lx+2.1, 1.95, 0.4, 0.6,
                     size=18, bold=True, color=col, align=PP_ALIGN.CENTER)

    add_divider(s, 4.38)

    # 图谱规模数据
    add_text(s, "图谱规模（满足作业要求）", 0.4, 4.5, 12, 0.45,
             size=14, bold=True, color=C_ACCENT)
    metrics = [
        ("110", "个实体", C_ACCENT),
        ("318", "条三元组", C_ACCENT2),
        ("  8", "类实体类型", C_GREEN),
        (" 16", "类关系类型", C_YELLOW),
        ("95.5%", "连通度", RGBColor(0x8E,0x44,0xAD)),
    ]
    for i, (val, lab, col) in enumerate(metrics):
        lx = 0.4 + i * 2.55
        add_rect(s, lx, 5.05, 2.1, 1.2, fill=C_CARD, line=col, line_w=Pt(2))
        add_text(s, val, lx, 5.1, 2.1, 0.65,
                 size=28, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(s, lab, lx, 5.72, 2.1, 0.45,
                 size=12, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_text(s, "※ 数据来源：芯片厂商官网 / AnandTech / IEEE Spectrum / 公开技术白皮书",
             0.4, 6.5, 12.5, 0.4, size=10, color=C_GRAY)


# ═══════════════════════════════════════════════════════════
# SLIDE 5 — 图谱样例与可视化
# ═══════════════════════════════════════════════════════════
def slide5(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "图谱样例与可视化", "典型三元组 · 中心性分析 · 交互式HTML图谱")

    # 典型三元组样例
    card(s, 0.3, 1.35, 6.1, 4.0, "典型三元组样例")
    triples = [
        ("(NVIDIA, designs, H100)",                    C_ACCENT),
        ("(TSMC, fabricates, H100)",                   C_ACCENT),
        ("(H100, uses_process, TSMC 4nm)",             C_ACCENT2),
        ("(H100, used_for, AI训练)",                    C_GREEN),
        ("(TSMC 3nm, successor_of, TSMC 4nm)",         C_YELLOW),
        ("(A17 Pro, uses_process, TSMC 3nm)",          C_ACCENT),
        ("(Kirin9000S, uses_process, SMIC N+1) [0.9]", C_ACCENT2),
        ("(EUV光刻, enables_process, TSMC 3nm)",        C_GREEN),
        ("(ASML, manufactures_tool, EUV光刻)",          C_ACCENT),
        ("(骁龙8Gen3, successor_of, 骁龙8Gen2)",        C_YELLOW),
    ]
    for i, (t, col) in enumerate(triples):
        add_rect(s, 0.4, 1.85 + i*0.33, 0.22, 0.22, fill=col)
        add_text(s, t, 0.7, 1.84 + i*0.33, 5.5, 0.32, size=11, color=C_WHITE)

    # 中心性分析
    card(s, 6.7, 1.35, 6.3, 2.4, "度中心性 Top-5（越高越核心）")
    central = [
        ("台积电",   0.239, C_ACCENT),
        ("英特尔",   0.184, C_ACCENT2),
        ("台积电3nm",0.174, C_GREEN),
        ("台积电4nm",0.156, C_YELLOW),
        ("骁龙8Gen3",0.138, RGBColor(0x8E,0x44,0xAD)),
    ]
    for i, (name, score, col) in enumerate(central):
        ty = 1.95 + i * 0.32
        bar_w = score * 16
        add_rect(s, 6.9, ty, bar_w, 0.24, fill=col)
        add_text(s, f"{name}  {score:.3f}", 6.95, ty+0.02, 6.0, 0.22,
                 size=11, bold=True, color=C_WHITE)

    # 可视化说明
    card(s, 6.7, 4.0, 6.3, 1.4, "交互式可视化（pyvis）")
    add_multiline(s, [
        ("results/kg_full.html  —  浏览器可拖拽、缩放、悬浮查看实体信息", False, C_WHITE),
        ("results/kg_nvidia_subgraph.html  —  NVIDIA 子图", False, C_LIGHT),
        ("results/kg_tsmc_subgraph.html    —  TSMC 子图", False, C_LIGHT),
    ], 6.85, 4.5, 6.0, 0.85, size=11)

    # 孤立节点说明
    card(s, 6.7, 5.6, 6.3, 1.45, "图谱质量说明")
    add_multiline(s, [
        ("最大连通分量：105/110（95.5%）", False, C_GREEN),
        ("孤立节点 5 个（部分材料实体待补充关系）", False, C_YELLOW),
        ("置信度字段显式标注争议信息", False, C_LIGHT),
    ], 6.85, 6.1, 6.0, 0.85, size=11)


# ═══════════════════════════════════════════════════════════
# SLIDE 6 — 推理方法
# ═══════════════════════════════════════════════════════════
def slide6(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "知识推理方法", "规则推理（Horn规则）+ 路径推理（双向BFS）")

    # 规则推理
    card(s, 0.3, 1.35, 6.1, 5.55, "方法一：规则推理（Horn 规则）")
    add_text(s, "规则形式：(A, rel1, B) ∧ (B, rel2, C)  ⟹  (A, new_rel, C)",
             0.45, 1.82, 5.8, 0.45, size=12, bold=True, color=C_ACCENT2)
    rules = [
        ("R1", "designs + uses_process → company_uses_process",   "0.95"),
        ("R2", "fabricates + uses_process → fab_masters_process", "0.90"),
        ("R3", "successor_of + competes_with → inherits_competition","0.85"),
        ("R5", "uses_device + uses_material → process_requires_material","0.90"),
        ("R6", "used_for + requires_device → chip_contains_device","0.75"),
    ]
    for i, (rn, desc, conf) in enumerate(rules):
        ty = 2.4 + i * 0.44
        add_rect(s, 0.4, ty, 0.42, 0.34, fill=C_ACCENT)
        add_text(s, rn, 0.4, ty, 0.42, 0.34, size=10, bold=True,
                 color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(s, desc, 0.9, ty+0.02, 4.4, 0.32, size=10, color=C_LIGHT)
        add_text(s, conf, 5.35, ty+0.02, 0.7, 0.32,
                 size=10, bold=True, color=C_GREEN)

    add_divider(s, 4.68, color=C_ACCENT2)
    add_text(s, "推理结果：7条规则共推导出 104 条新三元组\n其中置信度 ≥ 0.7 者写回图谱（inferred=True）",
             0.4, 4.75, 5.8, 0.75, size=12, color=C_YELLOW)

    add_text(s, "推理样例：",
             0.4, 5.6, 5.8, 0.35, size=12, bold=True, color=C_ACCENT)
    add_rect(s, 0.3, 5.95, 6.1, 0.8, fill=C_CARD, line=C_ACCENT, line_w=Pt(1))
    add_text(s, "已知：(TSMC, fabricates, H100) ∧ (H100, uses_process, TSMC 4nm)\n推导：(TSMC, fab_masters_process, TSMC 4nm)  [置信度: 0.90]",
             0.45, 6.0, 5.8, 0.72, size=11, color=C_WHITE)

    # 路径推理
    card(s, 6.7, 1.35, 6.3, 5.55, "方法二：路径推理（双向 BFS）")
    add_text(s, "在有向图上同时走正向边与反向边（inv_前缀标注）\n最大跳数 4，最多返回 5 条路径",
             6.85, 1.82, 6.0, 0.65, size=12, color=C_LIGHT)

    add_text(s, "示例：AMD → EUV光刻  (3跳路径)",
             6.85, 2.55, 6.0, 0.38, size=12, bold=True, color=C_ACCENT)
    add_rect(s, 6.7, 2.95, 6.3, 1.05, fill=C_CARD, line=C_ACCENT, line_w=Pt(1))
    add_text(s,
             "路径1：[AMD] --designs--> [Ryzen9 7950X]\n"
             "        --uses_process--> [TSMC 4nm]\n"
             "        --inv_enables_process--> [EUV光刻]",
             6.85, 2.98, 6.1, 0.98, size=10, color=C_WHITE)

    add_divider(s, 4.15, color=C_GREEN)
    add_text(s, "路径推理准确率评估（7个测试用例）",
             6.85, 4.22, 6.0, 0.38, size=12, bold=True, color=C_GREEN)

    cases = [
        ("AMD → EUV光刻（3跳）", "5条路径", True),
        ("苹果 → 台积电",         "5条路径", True),
        ("H100 → 硅材料",        "5条路径", True),
        ("Intel → AI训练应用",   "5条路径", True),
        ("麒麟9000S → SMIC工艺", "5条路径", True),
        ("NVIDIA → EUV技术",     "5条路径", True),
        ("Intel ≠ AMD直接路径",  "误找竞争路径", False),
    ]
    for i, (q, res, ok) in enumerate(cases):
        col = C_GREEN if ok else C_RED
        mark = "✓" if ok else "✗"
        add_text(s, f"{mark}  {q}  →  {res}",
                 6.85, 4.65 + i*0.32, 6.0, 0.3, size=11, color=col)

    add_text(s, "准确率：85.7%（6/7）", 6.85, 7.0, 6.0, 0.38,
             size=13, bold=True, color=C_GREEN)


# ═══════════════════════════════════════════════════════════
# SLIDE 7 — 大模型增强方法（GraphRAG）
# ═══════════════════════════════════════════════════════════
def slide7(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "大模型增强方法：GraphRAG", "图谱检索 → Prompt注入 → 对比问答")

    # 流程图（横向）
    steps7 = [
        ("①\n自然语言\n问题", C_ACCENT),
        ("②\n实体\n检索", C_ACCENT2),
        ("③\n子图\n扩展", C_GREEN),
        ("④\n三元组\n转文本", RGBColor(0x8E,0x44,0xAD)),
        ("⑤\nPrompt\n注入", C_YELLOW),
        ("⑥\nLLM\n生成", C_RED),
    ]
    for i, (label, col) in enumerate(steps7):
        lx = 0.3 + i * 2.15
        add_rect(s, lx, 1.4, 1.85, 1.2, fill=col)
        add_text(s, label, lx, 1.4, 1.85, 1.2,
                 size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        if i < 5:
            add_text(s, "▶", lx+1.85, 1.82, 0.28, 0.38,
                     size=16, bold=True, color=col, align=PP_ALIGN.CENTER)

    add_divider(s, 2.75)

    # 左列：Prompt模板约束
    card(s, 0.3, 2.92, 5.9, 3.55, "Prompt 设计约束（kg_augmented_prompt.txt）")
    add_multiline(s, [
        ("① 优先使用KG事实，用【KG事实】标记", True, C_ACCENT),
        ("② 无图谱证据时说明：未找到直接证据", True, C_YELLOW),
        ("③ 输出：直接答案 → KG三元组证据 → 置信度", True, C_GREEN),
        ("④ 多跳推理须显式列出推理链", True, C_ACCENT2),
    ], 0.5, 3.45, 5.6, 2.3, size=12)

    add_text(s, "{kg_context}  +  {question}", 0.45, 5.7, 5.5, 0.5,
             size=13, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_rect(s, 0.35, 5.65, 5.8, 0.5, fill=None, line=C_ACCENT, line_w=Pt(1))

    # 右列：幻觉检测 + 对比维度
    card(s, 6.5, 2.92, 6.5, 1.75, "幻觉检测（evaluate.py）")
    add_text(s, "给定LLM声明 → 在图谱中查询对应三元组\n存在 → 事实正确  ✅     不存在 → 可能幻觉  ❌",
             6.65, 3.42, 6.2, 0.9, size=12, color=C_WHITE)

    card(s, 6.5, 4.82, 6.5, 1.65, "四维对比指标")
    add_multiline(s, [
        ("准确性 — 工艺型号是否精确", False, C_WHITE),
        ("可追溯性 — 是否附有KG三元组列表", False, C_WHITE),
        ("幻觉抑制 — 错误声明是否被检出", False, C_WHITE),
        ("多跳推理 — 是否给出完整推理链", False, C_WHITE),
    ], 6.65, 5.1, 6.2, 1.25, size=12)

    add_text(s, "支持 DeepSeek API（真实问答）和规则模拟模式（离线演示）",
             0.3, 6.9, 12.7, 0.4, size=11, color=C_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 8 — 实验结果
# ═══════════════════════════════════════════════════════════
def slide8(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "实验结果", "检索评估 · 推理准确率 · 幻觉检测 · 对比分析")

    # 指标摘要行
    kpis = [
        ("59.0%",  "检索宏平均\n精确率",   C_ACCENT),
        ("65.8%",  "检索宏平均\n召回率",   C_ACCENT2),
        ("85.7%",  "路径推理\n准确率",     C_GREEN),
        ("100%",   "幻觉检测\n准确率",     C_YELLOW),
        ("104条",  "规则推理\n新三元组",   RGBColor(0x8E,0x44,0xAD)),
    ]
    for i, (val, lab, col) in enumerate(kpis):
        lx = 0.3 + i * 2.58
        add_rect(s, lx, 1.35, 2.35, 1.05, fill=col)
        add_text(s, val, lx, 1.35, 2.35, 0.62,
                 size=26, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(s, lab, lx, 1.92, 2.35, 0.45,
                 size=10, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_divider(s, 2.55)

    # 左下：逐题检索结果
    card(s, 0.3, 2.7, 7.5, 4.4, "逐题检索 Precision / Recall / F1")
    rows = [
        ("Q1  H100工艺（代工厂）",  "0.60", "0.75", "0.667", C_GREEN),
        ("Q2  台积电工艺（一对多）", "0.80", "1.00", "0.889", C_GREEN),
        ("Q3  A17 Pro vs 骁龙（比较）","0.40","0.67","0.500", C_YELLOW),
        ("Q4  麒麟9000S（争议）",   "1.00", "0.67", "0.800", C_GREEN),
        ("Q5  EUV供应商",           "0.40", "0.50", "0.444", C_YELLOW),
        ("Q6  AI训练芯片",          "0.60", "0.75", "0.667", C_GREEN),
        ("Q9  骁龙代际演进",        "1.00", "1.00", "1.000", C_GREEN),
        ("Q10 3nm材料（最低分）",   "0.20", "0.25", "0.222", C_RED),
    ]
    headers = ["问题", "P", "R", "F1"]
    for j, h in enumerate(headers):
        lx = [0.45, 5.1, 5.75, 6.5][j]
        add_text(s, h, lx, 3.1, [4.5,0.6,0.6,0.65][j], 0.3,
                 size=11, bold=True, color=C_ACCENT)
    for i, (q, p, r, f, col) in enumerate(rows):
        ty = 3.5 + i * 0.38
        add_text(s, q,  0.45, ty, 4.55, 0.34, size=10, color=C_LIGHT)
        add_text(s, p,  5.1,  ty, 0.55, 0.34, size=10, bold=True, color=col)
        add_text(s, r,  5.75, ty, 0.55, 0.34, size=10, bold=True, color=col)
        add_text(s, f,  6.5,  ty, 0.65, 0.34, size=10, bold=True, color=col)

    # 右下：幻觉检测明细
    card(s, 8.1, 2.7, 4.9, 4.4, "幻觉检测明细（8/8）")
    hall = [
        ("H100 采用 TSMC 4nm",       "✅ 正确"),
        ("M3 Ultra 用 Samsung 3nm",  "❌ 幻觉！"),
        ("Ryzen 7950X 基于 Zen4",    "✅ 正确"),
        ("Intel 设计了骁龙8Gen3",     "❌ 幻觉！"),
        ("TSMC 代工 A17 Pro",        "✅ 正确"),
        ("麒麟9000S 是麒麟9000继代", "✅ 正确"),
        ("ASML 是EUV唯一供应商",     "✅ 正确"),
        ("Ryzen 由 Intel 代工",      "❌ 幻觉！"),
    ]
    for i, (claim, result) in enumerate(hall):
        col = C_GREEN if "✅" in result else C_RED
        add_text(s, f"{result}  {claim}", 8.2, 3.1 + i*0.44, 4.65, 0.38,
                 size=10, color=col)


# ═══════════════════════════════════════════════════════════
# SLIDE 9 — 无KG vs 有KG 对比案例（系统演示）
# ═══════════════════════════════════════════════════════════
def slide9(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "案例对比演示", "无KG基线 vs GraphRAG增强回答")

    case_items = [
        {
            "q":    "Q1：苹果A17 Pro 和骁龙8 Gen 3，谁的制造工艺更先进？",
            "base": "A17 Pro可能采用3nm工艺，骁龙8 Gen 3也在4nm左右。（模糊，无确切证据）",
            "kg":   "【KG事实】A17 Pro→TSMC 3nm ✅  /  骁龙8Gen3→TSMC 4nm(N4P) ✅\n(TSMC 3nm, successor_of, TSMC 4nm) ✅\n→ A17 Pro（3nm）更先进；置信度：高",
        },
        {
            "q":    "Q5：EUV光刻技术的唯一供应商是哪家公司？",
            "base": "ASML是EUV光刻机的供应商，对先进制程至关重要。（正确但无证据）",
            "kg":   "【KG事实】(ASML, manufactures_tool, EUV光刻) ✅\n(EUV光刻, enables_process, TSMC 3nm/4nm/5nm) ✅\n→ ASML垄断EUV供应；具体受益工艺节点有据可查",
        },
        {
            "q":    "幻觉案例：LLM声明 [Apple M3 Ultra采用Samsung 3nm GAA工艺]",
            "base": "模型输出：M3 Ultra采用三星3nm工艺……（❌ 幻觉）",
            "kg":   "图谱核查：(M3 Ultra, uses_process, Samsung 3nm GAA) → 不存在\n正确事实：(M3 Ultra, uses_process, TSMC 3nm(N3)) ✅\n→ 幻觉检出！Apple芯片历来由TSMC代工",
        },
    ]

    for i, c in enumerate(case_items):
        ty = 1.4 + i * 1.96
        add_rect(s, 0.3, ty, 12.7, 0.35, fill=C_ACCENT)
        add_text(s, c["q"], 0.45, ty+0.04, 12.4, 0.3,
                 size=12, bold=True, color=C_WHITE)
        add_rect(s, 0.3, ty+0.35, 6.1, 1.4, fill=C_CARD, line=C_RED, line_w=Pt(1))
        add_text(s, "无KG基线", 0.35, ty+0.38, 2.0, 0.3,
                 size=10, bold=True, color=C_RED)
        add_text(s, c["base"], 0.45, ty+0.68, 5.85, 0.9, size=10, color=C_LIGHT)
        add_rect(s, 6.6, ty+0.35, 6.4, 1.4, fill=C_CARD, line=C_GREEN, line_w=Pt(1))
        add_text(s, "KG增强回答", 6.65, ty+0.38, 2.5, 0.3,
                 size=10, bold=True, color=C_GREEN)
        add_text(s, c["kg"], 6.7, ty+0.68, 6.1, 0.9, size=10, color=C_LIGHT)

    add_text(s, "交互演示：打开 results/kg_full.html，在图中追踪以上三元组", 0.3, 7.1, 12.7, 0.38,
             size=11, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 10 — 总结与不足
# ═══════════════════════════════════════════════════════════
def slide10(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "总结与不足", "工作总结 · 局限分析 · 未来方向")

    # 左：工作总结
    card(s, 0.3, 1.35, 6.1, 5.5, "主要工作与成果")
    add_multiline(s, [
        ("✅  图谱构建", True, C_ACCENT),
        ("  110实体、318三元组、8类实体、16类关系", False, C_WHITE),
        ("  连通率95.5%，pyvis可交互可视化", False, C_LIGHT),
        ("", False, C_WHITE),
        ("✅  知识推理", True, C_ACCENT2),
        ("  7条Horn规则，推导104条新三元组", False, C_WHITE),
        ("  双向BFS路径推理，准确率85.7%", False, C_LIGHT),
        ("", False, C_WHITE),
        ("✅  大模型增强（GraphRAG）", True, C_GREEN),
        ("  图谱子图检索 → Prompt注入 → 对比问答", False, C_WHITE),
        ("  幻觉检测准确率100%（8/8声明）", False, C_LIGHT),
        ("  检索宏P=59%，宏R=65.8%", False, C_LIGHT),
    ], 0.5, 1.88, 5.7, 4.8, size=12)

    # 右上：局限性
    card(s, 6.7, 1.35, 6.3, 3.2, "局限性")
    add_multiline(s, [
        ("⚠  图谱覆盖不完整，5个孤立节点", False, C_YELLOW),
        ("⚠  材料类实体跨类型检索召回率低（Q10: F1=0.22）", False, C_YELLOW),
        ("⚠  竞争继承规则仅推导1条，关系链稀疏", False, C_YELLOW),
        ("⚠  麒麟9000S代工厂信息存在争议标注不一致", False, C_YELLOW),
        ("⚠  未实现TransE等表示学习推理方法", False, C_YELLOW),
    ], 6.85, 1.85, 6.0, 2.5, size=12)

    # 右下：未来方向
    card(s, 6.7, 4.7, 6.3, 2.15, "未来改进方向")
    add_multiline(s, [
        ("→  引入语义相似度提升跨类型实体检索", False, C_GREEN),
        ("→  半自动抽取（LLM生成候选+人工审核）扩充图谱", False, C_GREEN),
        ("→  实现TransE/RotatE提升链接预测效果", False, C_GREEN),
        ("→  Text2Cypher支持结构化自然语言查询", False, C_GREEN),
    ], 6.85, 5.18, 6.0, 1.55, size=12)

    # 底部致谢
    add_rect(s, 0, 7.1, 13.33, 0.4, fill=C_ACCENT)
    add_text(s, "感谢老师和同学的评阅与指正！", 0, 7.12, 13.33, 0.35,
             size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 11 — 小组分工
# ═══════════════════════════════════════════════════════════
def slide11(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "小组分工", "三位成员各司其职，共同完成完整系统")

    members = [
        {
            "name": "成员 A",
            "color": C_ACCENT,
            "tasks": [
                "数据收集与整理",
                "查阅芯片厂商官网、行业媒体、技术白皮书",
                "设计 8 类实体 / 16 类关系 Schema",
                "完成 entities.csv 与 triples.csv 建立核对",
                "争议信息置信度标注（如麒麟9000S）",
                "参与第二章报告撰写",
            ],
        },
        {
            "name": "成员 B",
            "color": C_ACCENT2,
            "tasks": [
                "核心代码实现",
                "build_kg.py — 图谱构建与NetworkX存储",
                "reasoning.py — 规则推理与路径推理",
                "visualize.py — pyvis交互式可视化",
                "图谱质量分析 / 路径推理实验",
                "参与第三章报告撰写",
            ],
        },
        {
            "name": "成员 C",
            "color": C_GREEN,
            "tasks": [
                "大模型增强模块",
                "graph_retrieval.py — 子图检索",
                "llm_qa.py — GraphRAG 问答对比",
                "evaluate.py — 综合评估与幻觉检测",
                "run_all.py — 一键运行脚本",
                "PPT 制作 / 第一、四、五章报告",
            ],
        },
    ]

    for i, m in enumerate(members):
        lx = 0.4 + i * 4.3
        add_rect(s, lx, 1.4, 4.0, 0.55, fill=m["color"])
        add_text(s, m["name"], lx, 1.4, 4.0, 0.55,
                 size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        for j, task in enumerate(m["tasks"]):
            ty = 2.08 + j * 0.72
            if j == 0:
                add_rect(s, lx, ty, 4.0, 0.6, fill=C_CARD, line=m["color"], line_w=Pt(2))
                add_text(s, task, lx+0.12, ty+0.1, 3.8, 0.42,
                         size=13, bold=True, color=m["color"])
            else:
                add_rect(s, lx, ty, 0.22, 0.22, fill=m["color"])
                add_text(s, task, lx+0.28, ty+0.01, 3.65, 0.38,
                         size=12, color=C_WHITE)

    add_divider(s, 6.72)
    add_text(s, "三位成员共同参与：系统设计讨论  ·  实验结果分析  ·  报告审阅  ·  答辩准备",
             0.4, 6.82, 12.5, 0.45, size=12, color=C_LIGHT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# 主程序
# ═══════════════════════════════════════════════════════════
def main():
    prs = new_prs()

    print("生成 Slide 1 — 封面...")
    slide1(prs)
    print("生成 Slide 2 — 研究背景与任务目标...")
    slide2(prs)
    print("生成 Slide 3 — Schema 设计...")
    slide3(prs)
    print("生成 Slide 4 — 数据来源与构建流程...")
    slide4(prs)
    print("生成 Slide 5 — 图谱样例与可视化...")
    slide5(prs)
    print("生成 Slide 6 — 推理方法...")
    slide6(prs)
    print("生成 Slide 7 — 大模型增强方法...")
    slide7(prs)
    print("生成 Slide 8 — 实验结果...")
    slide8(prs)
    print("生成 Slide 9 — 系统演示案例...")
    slide9(prs)
    print("生成 Slide 10 — 总结与不足...")
    slide10(prs)
    print("生成 Slide 11 — 小组分工...")
    slide11(prs)

    out = os.path.join(os.path.dirname(__file__), "课程PPT.pptx")
    prs.save(out)
    print(f"\n✅ PPT已保存：{out}")
    print(f"   共 {len(prs.slides)} 张幻灯片")


if __name__ == "__main__":
    main()
